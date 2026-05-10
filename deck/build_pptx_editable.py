#!/usr/bin/env python3
"""
Build a fully EDITABLE Phetaloun SchoolBus deck (PowerPoint widescreen 16:9).

Every heading, body text, card, pill, and divider is a native PowerPoint
shape you can click and edit. Only the product screenshots are embedded
as images.

Run: python3 build_pptx_editable.py
Output: Phetaloun-SchoolBus-Lao-editable.pptx (next to this script)
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------------------------------------------------
DECK_DIR  = Path(__file__).resolve().parent
SHOTS_DIR = DECK_DIR / "screenshots"
OUT       = DECK_DIR / "Phetaloun-SchoolBus-Lao-editable.pptx"

# PowerPoint widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Phetaloun palette
BLUE      = RGBColor(0x0E, 0x4D, 0x92)
BLUE_2    = RGBColor(0x13, 0x4A, 0xA6)
CYAN      = RGBColor(0x29, 0xAB, 0xE2)
CYAN_DARK = RGBColor(0x1B, 0x95, 0xC9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT      = RGBColor(0x0F, 0x17, 0x2A)
TEXT_2    = RGBColor(0x33, 0x41, 0x55)
TEXT_3    = RGBColor(0x64, 0x74, 0x8B)
TEXT_4    = RGBColor(0x94, 0xA3, 0xB8)
BG        = RGBColor(0xF1, 0xF5, 0xF9)
MUTED     = RGBColor(0xF8, 0xFA, 0xFC)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
RED       = RGBColor(0xEF, 0x44, 0x44)
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)

# Recommended Lao font that also handles Latin
# Note: Lao MN (macOS) and Phetsarath OT (Windows/Linux) both render Lao correctly.
# We set Lao MN since it's the macOS system default; PowerPoint on Windows
# will substitute Phetsarath OT or another Lao font automatically.
FONT_LAO  = "Lao MN"
FONT_LAT  = "Calibri"  # fallback for English-only labels

# ----------------------------------------------------------------------
def set_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_fill(shape):
    shape.fill.background()

def no_line(shape):
    shape.line.fill.background()

def thin_line(shape, color, width_pt=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)

def set_para_font(p, name=FONT_LAO, size=14, bold=False, color=TEXT, align=None):
    if align is not None:
        p.alignment = align
    for r in p.runs:
        r.font.name = name
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

def add_text(slide, x, y, w, h, text,
             size=14, bold=False, color=TEXT, name=FONT_LAO,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """Add a textbox; supports multi-paragraph via \n."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    paras = str(text).split('\n')
    for i, line in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for r in p.runs:
            r.font.name = name
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return tb

def add_pill(slide, x, y, text, fill=CYAN, color=WHITE, size=11):
    """Small rounded label pill."""
    width  = Inches(max(0.6, 0.16 * len(text) + 0.4))
    height = Inches(0.32)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    shape.adjustments[0] = 0.5
    set_solid(shape, fill)
    no_line(shape)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Emu(60000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.name = FONT_LAO
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = color
    return shape

def add_card(slide, x, y, w, h, fill=WHITE, border=BORDER, corner=0.04):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = corner
    set_solid(shape, fill)
    if border is None:
        no_line(shape)
    else:
        thin_line(shape, border, 1)
    no_text_frame_default(shape)
    return shape

def no_text_frame_default(shape):
    # remove the default empty paragraph styling so it doesn't add visible space
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

def add_rect(slide, x, y, w, h, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_solid(shape, fill)
    if border is None:
        no_line(shape)
    else:
        thin_line(shape, border, 1)
    return shape

def add_circle(slide, x, y, d, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    set_solid(shape, fill)
    if border is None:
        no_line(shape)
    else:
        thin_line(shape, border, 2)
    return shape

def add_image(slide, image_path, x, y, w=None, h=None):
    if not Path(image_path).exists():
        print(f"  [warn] missing image: {image_path}", file=sys.stderr)
        return None
    if w and h:
        return slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(str(image_path), x, y, width=w)
    return slide.shapes.add_picture(str(image_path), x, y, height=h)

def add_slide_number(slide, n, total, dark=False):
    tb = add_text(slide,
        Inches(12.4), Inches(7.10), Inches(0.8), Inches(0.3),
        f"{n} / {total}",
        size=10, bold=False,
        color=RGBColor(0xFF,0xFF,0xFF) if dark else TEXT_4,
        align=PP_ALIGN.RIGHT,
        name=FONT_LAT)
    return tb

# ----------------------------------------------------------------------
# Slide builders

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

TOTAL = 13

def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    bg_rect = add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=bg)
    bg_rect.shadow.inherit = False
    return s

def new_brand_slide():
    s = prs.slides.add_slide(BLANK)
    # Solid blue background (PowerPoint widescreen full-bleed)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=BLUE)
    return s

# ---------- Slide 1: TITLE ----------
def slide_1():
    s = new_brand_slide()
    # Logo block (white card with logo image)
    logo_card = add_rect(s, Inches(0.7), Inches(2.4), Inches(2.6), Inches(1.0), fill=WHITE)
    logo_card.adjustments
    logo_path = DECK_DIR.parent / "shared" / "phetaloun-logo.svg"
    # python-pptx doesn't reliably embed SVG, so use PNG fallback if available
    logo_png = DECK_DIR / "phetaloun-logo.png"
    if logo_png.exists():
        add_image(s, logo_png, Inches(0.85), Inches(2.55), w=Inches(2.3))
    else:
        # Fallback: text-based brand
        add_text(s, Inches(0.85), Inches(2.65), Inches(2.3), Inches(0.6),
            "Phetaloun", size=28, bold=True, color=BLUE, name=FONT_LAT)

    # Hero title
    add_text(s, Inches(0.7), Inches(3.6), Inches(12), Inches(1.4),
        "Phetaloun SchoolBus", size=64, bold=True, color=WHITE, name=FONT_LAT)

    # Tagline (Lao)
    add_text(s, Inches(0.7), Inches(4.9), Inches(12), Inches(1.0),
        "ລະບົບປະຕິບັດການລົດເມໂຮງຮຽນສຳລັບແຂວງບໍ່ແກ້ວ. ຮັບສົ່ງເຖິງປະຕູບ້ານ, ສະແກນ NFC + QR,\nຕິດຕາມສົດ, ການຈັດການຄະແນນ ແລະ ການຂາດຮຽນ — ຄົບໃນລະບົບດຽວ.",
        size=18, color=RGBColor(0xE0,0xEC,0xFA), line_spacing=1.4)

    # Pills
    add_pill(s, Inches(0.7), Inches(6.2), "ຕົວແບບ · 2026", fill=RGBColor(0xFF,0xFF,0xFF), color=BLUE, size=11)
    add_pill(s, Inches(2.4), Inches(6.2), "ຫ້ວຍຊາຍ · ບໍ່ແກ້ວ", fill=CYAN, color=WHITE, size=11)

    add_slide_number(s, 1, TOTAL, dark=True)

# ---------- Slide 2: WHAT IS ----------
def slide_2():
    s = new_slide()
    add_pill(s, Inches(0.7), Inches(0.7), "01 · ພາບລວມ", fill=BG, color=BLUE)
    add_text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.8),
        "Phetaloun SchoolBus ແມ່ນຫຍັງ?", size=32, bold=True, color=BLUE)
    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.4),
        "ລະບົບຄົບຖ້ວນສຳລັບບໍລິສັດທີ່ໃຫ້ບໍລິການລົດເມໂຮງຮຽນຫຼາຍແຫ່ງ. ລົງທະບຽນນັກຮຽນ, ປັກຕຳແໜ່ງເຮືອນ, ວາງແຜນເສັ້ນທາງຮັບສົ່ງເຖິງປະຕູ, ສະແກນບັດທີ່ປະຕູລົດ, ຕິດຕາມສົດ ແລະ ແຈ້ງເຕືອນຜູ້ປົກຄອງທັນທີ.",
        size=15, color=TEXT_2, line_spacing=1.5)

    cards = [
        ("ບໍລິສັດດຽວ ຫຼາຍໂຮງຮຽນ",
         "ໜ້າຄວບຄຸມດຽວສຳລັບກອງລົດທັງໝົດ. ຈັດການໂຮງຮຽນ, ລົດ, ຄົນຂັບ, ເສັ້ນທາງ, ບັນຊີຜູ້ໃຊ້ລະບົບ.",
         CYAN),
        ("ຮັບສົ່ງເຖິງປະຕູບ້ານ",
         "ປັກໝຸດເຮືອນຂອງແຕ່ລະນັກຮຽນ. ປັບປຸງເສັ້ນທາງໃຫ້ສັ້ນສຸດ. ຮັບເຖິງໜ້າປະຕູບ່ໍຕ້ອງຍ່າງໄປປ້າຍຈອດ.",
         AMBER),
        ("ສະແກນດ້ວຍ NFC + QR",
         "ບັດນັກຮຽນມີທັງ NFC ແລະ QR. ຄົນຂັບໃຊ້ PDA ມືຖືແຕະ ຫຼື ສະແກນ. ປະຫຍັດເວລາ ຫຼຸດຄວາມຜິດພາດ.",
         GREEN),
    ]
    cw = Inches(3.95)
    cgap = Inches(0.18)
    cx = Inches(0.7)
    cy = Inches(4.0)
    ch = Inches(2.7)
    for i, (title, body, accent) in enumerate(cards):
        x = cx + (cw + cgap) * i
        c = add_card(s, x, cy, cw, ch)
        # accent bar at top
        add_rect(s, x, cy, cw, Inches(0.12), fill=accent)
        add_text(s, x + Inches(0.3), cy + Inches(0.4), cw - Inches(0.6), Inches(0.6),
            title, size=18, bold=True, color=BLUE)
        add_text(s, x + Inches(0.3), cy + Inches(1.05), cw - Inches(0.6), Inches(1.5),
            body, size=13, color=TEXT_2, line_spacing=1.5)
    add_slide_number(s, 2, TOTAL)

# ---------- Slide 3: PROBLEM ----------
def slide_3():
    s = new_slide()
    add_pill(s, Inches(0.7), Inches(0.7), "02 · ບັນຫາ", fill=BG, color=BLUE)
    add_text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.8),
        "ບັນຫາທີ່ພົບໃນປະຈຸບັນ", size=32, bold=True, color=BLUE)
    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.0),
        "ປະຈຸບັນ ຜູ້ປົກຄອງ ໂຮງຮຽນ ແລະ ບໍລິສັດລົດເມ ບໍ່ມີຂໍ້ມູນທີ່ຊັດເຈນ ເຮັດໃຫ້ບໍ່ສາມາດຕິດຕາມການເຄື່ອນໄຫວ ແລະ ຄວາມຜິດພາດເກີດຂຶ້ນເລື້ອຍໆ.",
        size=15, color=TEXT_2, line_spacing=1.5)

    items = [
        ("ຜູ້ປົກຄອງເບິ່ງບໍ່ເຫັນ",  "ໂທໂຮງຮຽນບໍ່ຮູ້ຈົບ ແຕ່ບໍ່ມີຄຳຕອບ.",                         RED),
        ("ຂຶ້ນຜິດລົດ",              "ໂດຍສະເພາະນັກຮຽນນ້ອຍ. ຜິດເທື່ອດຽວ ກໍ່ເປັນຂ່າວ.",            AMBER),
        ("ກວດເຂົ້າຮຽນດ້ວຍເຈ້ຍ",   "ບັນທຶກດ້ວຍມື. ປຽບທຽບທາງໂທລະສັບ. ເສຍເວລາ.",                AMBER),
        ("ເດັກຄ້າງໃນລົດ",            "ເຫດການຈິງທີ່ເກີດຊ້ຳໆ. ເສຍຊື່ສຽງຍາວນານ.",                    RED),
    ]
    cw = Inches(2.95)
    cgap = Inches(0.15)
    cx = Inches(0.7)
    cy = Inches(3.6)
    ch = Inches(3.0)
    for i, (title, body, accent) in enumerate(items):
        x = cx + (cw + cgap) * i
        add_card(s, x, cy, cw, ch)
        # icon circle
        add_circle(s, x + Inches(0.3), cy + Inches(0.35), Inches(0.55), fill=accent)
        add_text(s, x + Inches(0.3), cy + Inches(1.1), cw - Inches(0.6), Inches(0.7),
            title, size=15, bold=True, color=BLUE)
        add_text(s, x + Inches(0.3), cy + Inches(1.85), cw - Inches(0.6), Inches(1.0),
            body, size=12, color=TEXT_2, line_spacing=1.5)
    add_slide_number(s, 3, TOTAL)

# ---------- Slide 4: 4 EVENTS ----------
def slide_4():
    s = new_slide()
    add_pill(s, Inches(0.7), Inches(0.7), "03 · ວິທີໃຊ້ງານ", fill=BG, color=BLUE)
    add_text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.9),
        "ສະແກນ 4 ຄັ້ງຕໍ່ມື້ — ໝົດເລື່ອງເທົ່ານັ້ນ", size=30, bold=True, color=BLUE)
    add_text(s, Inches(0.7), Inches(2.1), Inches(12), Inches(1.0),
        "ທຸກມື້ໂຮງຮຽນ ທຸກນັກຮຽນ. ການສະແກນແຕ່ລະຄັ້ງສ້າງເຫດການທີ່ກວດສອບໄດ້ ແລະ ສົ່ງການແຈ້ງເຕືອນຫາຜູ້ປົກຄອງທັນທີ.",
        size=15, color=TEXT_2, line_spacing=1.5)

    # 4 event cards in a row with arrows
    events = [
        ("ເຫດການ 1", "ຂຶ້ນລົດທີ່ບ້ານ",   "ເຊົ້າ"),
        ("ເຫດການ 2", "ມາຮອດໂຮງຮຽນ",     "ເຊົ້າ"),
        ("ເຫດການ 3", "ຂຶ້ນລົດກັບບ້ານ",   "ແລງ"),
        ("ເຫດການ 4", "ສົ່ງເຖິງບ້ານ",         "ແລງ"),
    ]
    cw = Inches(2.55)
    cgap_x = Inches(0.25)
    cy = Inches(3.6)
    ch = Inches(2.4)
    cx = Inches(0.85)
    for i, (num, title, when) in enumerate(events):
        x = cx + (cw + cgap_x) * i
        c = add_card(s, x, cy, cw, ch, border=BORDER)
        # arrow between cards (except after last)
        if i < len(events) - 1:
            arrow_x = x + cw + Inches(0.02)
            arrow_y = cy + ch / 2 - Inches(0.15)
            arrow = slide_shape_arrow(s, arrow_x, arrow_y, Inches(0.21), Inches(0.3))
        add_text(s, x, cy + Inches(0.25), cw, Inches(0.3),
            num, size=10, bold=True, color=TEXT_3, align=PP_ALIGN.CENTER)
        # icon (use cyan circle)
        add_circle(s, x + cw/2 - Inches(0.35), cy + Inches(0.7), Inches(0.7), fill=CYAN)
        add_text(s, x, cy + Inches(1.55), cw, Inches(0.6),
            title, size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(2.05), cw, Inches(0.3),
            when, size=12, color=TEXT_3, align=PP_ALIGN.CENTER)

    # Footer note
    note_y = Inches(6.4)
    note_card = add_card(s, Inches(0.7), note_y, Inches(11.93), Inches(0.65), fill=BG, border=None)
    # left accent
    add_rect(s, Inches(0.7), note_y, Inches(0.08), Inches(0.65), fill=CYAN)
    add_text(s, Inches(0.95), note_y + Inches(0.13), Inches(11.5), Inches(0.4),
        "ແຕ່ລະເຫດການບັນທຶກ: ນັກຮຽນ · ລົດ · ປ້າຍຈອດ · ເວລາ · GPS · ຄົນຂັບ.",
        size=13, color=TEXT_2)

    add_slide_number(s, 4, TOTAL)

def slide_shape_arrow(s, x, y, w, h):
    # right-arrow
    shape = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    set_solid(shape, CYAN)
    no_line(shape)
    return shape

# ---------- Slide 5–9 + 12: image-left/right with bullets ----------
def slide_image_bullets(idx, pill_text, title, intro, bullets, image_name, image_side="left"):
    s = new_slide()
    add_pill(s, Inches(0.7), Inches(0.7), pill_text, fill=BG, color=BLUE)
    add_text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.9),
        title, size=28, bold=True, color=BLUE)

    # Layout: image takes ~7" wide, bullets ~5" wide
    img_w = Inches(7.0)
    img_x = Inches(0.7) if image_side == "left" else Inches(5.6)
    img_y = Inches(2.3)
    add_image(s, SHOTS_DIR / f"{image_name}.png", img_x, img_y, w=img_w)

    txt_x = Inches(8.0) if image_side == "left" else Inches(0.7)
    txt_y = Inches(2.3)
    txt_w = Inches(4.7)
    add_text(s, txt_x, txt_y, txt_w, Inches(0.7),
        intro, size=14, color=TEXT_2, line_spacing=1.5)

    by = txt_y + Inches(1.2)
    for b in bullets:
        # Tick + text row
        add_text(s, txt_x, by, Inches(0.3), Inches(0.36), "✓",
            size=14, bold=True, color=CYAN, name=FONT_LAT)
        add_text(s, txt_x + Inches(0.38), by, txt_w - Inches(0.4), Inches(0.7),
            b, size=13, color=TEXT_2, line_spacing=1.4)
        by += Inches(0.6)

    add_slide_number(s, idx, TOTAL)

# ---------- Slide 10–11: side-by-side flow comparison ----------
def slide_two_screens(idx, pill_text, title, left_label, left_img, right_label, right_img, footnote):
    s = new_slide()
    add_pill(s, Inches(0.7), Inches(0.7), pill_text, fill=RGBColor(0xFE,0xF3,0xC7), color=AMBER)
    add_text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.9),
        title, size=26, bold=True, color=BLUE)

    iw = Inches(5.5)
    ix1 = Inches(0.85)
    ix2 = Inches(7.0)
    iy = Inches(2.5)

    add_text(s, ix1, iy - Inches(0.35), iw, Inches(0.32),
        left_label, size=14, bold=True, color=BLUE)
    add_text(s, ix2, iy - Inches(0.35), iw, Inches(0.32),
        right_label, size=14, bold=True, color=BLUE)

    add_image(s, SHOTS_DIR / f"{left_img}.png",  ix1, iy, w=iw)
    add_image(s, SHOTS_DIR / f"{right_img}.png", ix2, iy, w=iw)

    # Footer note card
    fy = Inches(6.45)
    add_card(s, Inches(0.7), fy, Inches(11.93), Inches(0.65), fill=BG, border=None)
    add_rect(s, Inches(0.7), fy, Inches(0.08), Inches(0.65), fill=AMBER)
    add_text(s, Inches(0.95), fy + Inches(0.12), Inches(11.5), Inches(0.45),
        footnote, size=12, color=TEXT_2, line_spacing=1.4)

    add_slide_number(s, idx, TOTAL)

# ---------- Slide 13: ASK ----------
def slide_13():
    s = new_brand_slide()
    add_pill(s, Inches(0.7), Inches(0.9), "12 · ຂັ້ນຕອນຕໍ່ໄປ", fill=RGBColor(0xFF,0xFF,0xFF), color=BLUE)
    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(1.4),
        "ມາທົດລອງນຳກັນ.", size=56, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(1.5),
        "ເລືອກ 2 ລົດ ແລະ 1 ໂຮງຮຽນ ສຳລັບການທົດລອງ 30 ວັນ. ພວກເຮົາຕິດຕັ້ງ PDA, ອອກບັດ NFC+QR ໃຫ້ນັກຮຽນ,\nຝຶກອົບຮົມຄົນຂັບ ແລະ ແນະນຳຜູ້ປົກຄອງ. ທ່ານນຳເອົາລົດ ແລະ ຄວາມສຳພັນກັບໂຮງຮຽນ.",
        size=18, color=RGBColor(0xE0,0xEC,0xFA), line_spacing=1.5)

    # 4 stat boxes
    stats = [
        ("30", "ວັນ"),
        ("2",  "ລົດ"),
        ("1",  "ໂຮງຮຽນ"),
        ("0",  "ອຸປະກອນທີ່ຕ້ອງຊື້ກ່ອນ"),
    ]
    sw = Inches(2.7)
    gap = Inches(0.2)
    sx = Inches(0.7)
    sy = Inches(5.05)
    for i, (num, lbl) in enumerate(stats):
        x = sx + (sw + gap) * i
        add_text(s, x, sy, sw, Inches(0.9), num,
            size=56, bold=True, color=CYAN, name=FONT_LAT)
        add_text(s, x, sy + Inches(0.95), sw, Inches(0.4), lbl,
            size=12, color=RGBColor(0xE0,0xEC,0xFA))

    # Contact line
    add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
        "Phetaloun SchoolBus  ·  ops@phetaloun.la  ·  +856 84 555 0100  ·  phoudssp.github.io/schoolbus",
        size=11, color=RGBColor(0xC7,0xD8,0xEE), name=FONT_LAT)

    add_slide_number(s, 13, TOTAL, dark=True)

# ----------------------------------------------------------------------
# Build all slides

slide_1()
slide_2()
slide_3()
slide_4()

slide_image_bullets(
    5, "04 · ບໍລິສັດລົດເມ", "ໜ້າຄວບຄຸມຜູ້ດູແລລະບົບສູງສຸດ",
    "ໜ້າຈໍດຽວ ບໍລິຫານທຸກຢ່າງ ທົ່ວທຸກໂຮງຮຽນທີ່ບໍລິການ.",
    [
        "ແຜນທີ່ສົດໆຂອງລົດທຸກຄັນ",
        "ສະຖິຕິລາຍວັນ ນັກຮຽນ ລາຍຮັບ ອັດຕາມາຕາມເວລາ",
        "ການແຈ້ງເຕືອນສົດໆ ຕາຕະລາງມື້ນີ້",
        "13 ໂມດູນ: ໂຮງຮຽນ, ນັກຮຽນ, ຜູ້ປົກຄອງ, ຄົນຂັບ, ລົດ, ເສັ້ນທາງ, ແພັກເກັດ, ການເງິນ, ລາຍງານ",
    ],
    "op-dashboard", image_side="left"
)

slide_image_bullets(
    6, "05 · ຕິດຕາມສົດ", "ຕິດຕາມສົດ + ຮັບສົ່ງເຖິງປະຕູ",
    "ລົດເຄື່ອນໄຫວຕາມເສັ້ນທາງຈິງ ຢຸດທີ່ປ້າຍຈອດ ແລະ ສົ່ງເຫດການເຂົ້າມາສົດໆ.",
    [
        "ລົດເຄື່ອນໄຫວສະໝຳ່ສະເໝີ ປັບຄວາມໄວຈຳລອງໄດ້",
        "ໝຸດເຮືອນຂອງນັກຮຽນທຸກຄົນເພື່ອເຫັນຖານລູກຄ້າ",
        "ລົ໋ອກໃນພື້ນທີ່ບໍ່ແກ້ວເທົ່ານັ້ນ — ບໍ່ຂ້າມໄທ",
        "ຄລິກລົດໃດກໍ່ໄດ້ ເບິ່ງລາຍລະອຽດ ຮອບປະຈຸບັນ",
    ],
    "op-monitoring", image_side="left"
)

slide_image_bullets(
    7, "06 · ເສັ້ນທາງ", "ວາງແຜນເສັ້ນທາງ ຮັບເຖິງໜ້າປະຕູ",
    "ປັກໝຸດເຮືອນຂອງແຕ່ລະນັກຮຽນ. ສະແດງປ້າຍຈອດ + ບ້ານ + ໂຮງຮຽນ ໃນແຜນທີ່ດຽວ.",
    [
        "ປຸ່ມ \"ປັບປຸງ\" ຫາເສັ້ນທາງສັ້ນສຸດ",
        "ສະແດງໄລຍະທາງ km ແລະ ເວລາໂດຍປະມານ",
        "ຍ້າຍປ້າຍຈອດຂຶ້ນລົງ ເພີ່ມ/ລຶບໄດ້",
        "ຮັບແຍກແຕ່ລະຫຼັງ ຫຼື ປ້າຍຈອດດຽວກັນ",
    ],
    "op-routes", image_side="left"
)

slide_image_bullets(
    8, "07 · ຄົນຂັບ", "ອຸປະກອນ PDA ມືຖືສຳລັບຄົນຂັບ",
    "PDA ມືຖືແຂງແຮງ ມີຊ່ອງສະແກນ QR/ບາໂຄດ + NFC ໃນຕົວ.",
    [
        "ຊ່ອງສະແກນ QR/ບາໂຄດ ດ້ານເທິງ ມີລາເຊີ",
        "ປຸ່ມສະແກນສີເຫຼືອງໃຫຍ່ ດ້ານລ່າງ",
        "ປຸ່ມ NFC ສີຟ້າ ສຳລັບແຕະບັດ",
        "ບັງຄັບກວດກາພາຍໃນລົດ ຫຼັງສຳເລັດເສັ້ນທາງ",
        "ລົງທະບຽນແບບກຳມື ສຳລັບລືມບັດ",
    ],
    "driver-pda", image_side="right"
)

slide_image_bullets(
    9, "08 · ຜູ້ປົກຄອງ", "ແອັບໂທລະສັບສຳລັບຜູ້ປົກຄອງ",
    "ຮັບແຈ້ງເຕືອນທຸກຄັ້ງທີ່ລູກຂຶ້ນ-ລົງລົດ. ເບິ່ງຕຳແໜ່ງສົດໆ ແລະ ເຮັດທຸກຢ່າງຈາກໂທລະສັບ.",
    [
        "ຕຳແໜ່ງລົດສົດໆເທິງແຜນທີ່ ພ້ອມເວລາມາຮອດ",
        "ແຈ້ງເຕືອນທຸກໆເຫດການທັງ 4 ຄັ້ງ",
        "ບັນຊີດຽວ ສຳລັບຫຼາຍລູກ ຄົນລະໂຮງຮຽນກໍ່ໄດ້",
        "ເບິ່ງຄະແນນ ແລະ ຂໍຂາດຮຽນຈາກໂທລະສັບ",
        "ສອງພາສາ ລາວ + ອັງກິດ",
    ],
    "parent-home", image_side="right"
)

slide_two_screens(
    10, "ໃໝ່ · 09 · ຄະແນນ", "ການຈັດການຄະແນນ — ໂຮງຮຽນ ↔ ຜູ້ປົກຄອງ",
    "ໂຮງຮຽນ: ໃສ່ຄະແນນ", "school-scores",
    "ຜູ້ປົກຄອງ: ເຫັນຄະແນນທັນທີ", "parent-scores",
    "ການເຮັດວຽກ: ຄູໃສ່ຄະແນນຢູ່ໂຮງຮຽນ → ຜູ້ປົກຄອງເຫັນຄະແນນ ແລະ ຄະແນນສະເລ່ຍຂອງລູກໃນແອັບໂທລະສັບທັນທີ. ມີກຣາບຄະແນນສະເລ່ຍຕາມວິຊາ."
)

slide_two_screens(
    11, "ໃໝ່ · 10 · ການຂາດຮຽນ", "ຂັ້ນຕອນຂໍຂາດຮຽນ — ຜູ້ປົກຄອງ → ໂຮງຮຽນ",
    "ຜູ້ປົກຄອງ: ຂໍຜ່ານແອັບ", "parent-absence",
    "ໂຮງຮຽນ: ອະນຸມັດ ຫຼື ປະຕິເສດ", "school-attendance",
    "ການເຮັດວຽກ: ຜູ້ປົກຄອງເລືອກລູກ → ວັນທີຂາດ → ເຫດຜົນ → ສົ່ງ. ໂຮງຮຽນເຫັນທັນທີ ກົດອະນຸມັດ ຫຼື ປະຕິເສດ."
)

slide_image_bullets(
    12, "11 · ໂຮງຮຽນ", "ໜ້າຄວບຄຸມໂຮງຮຽນ + ຈັດການວິຊາ/ປະເພດ/ຄູ",
    "ແຕ່ລະໂຮງຮຽນເຫັນຂໍ້ມູນສະເພາະຂອງຕົນ ແລະ ຈັດການລະບົບການສຶກສາໄດ້ດ້ວຍຕົນເອງ.",
    [
        "ສະຖານະການລົງທະບຽນ ການເຂົ້າຮຽນ ສົດໆ",
        "CRUD ວິຊາຮຽນ ເພີ່ມ/ແກ້ໄຂ/ລຶບ",
        "CRUD ປະເພດການປະເມີນ ມີນ້ຳໜັກ %",
        "CRUD ຄູອາຈານ ພ້ອມວິຊາທີ່ສອນ",
        "ແບບຟອມໃສ່ຄະແນນ ຮັບລາຍຊື່ຈາກ CRUD",
    ],
    "school-dashboard", image_side="left"
)

slide_13()

# ----------------------------------------------------------------------
prs.save(OUT)
print(f"Wrote {OUT} ({len(prs.slides)} slides, fully editable)")
