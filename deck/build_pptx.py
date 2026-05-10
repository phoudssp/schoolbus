#!/usr/bin/env python3
"""Assemble per-slide PNGs into a PowerPoint widescreen 16:9 deck."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

DECK_DIR = Path(__file__).resolve().parent
SLIDES_DIR = DECK_DIR / "slides"
OUT = DECK_DIR / "Phetaloun-SchoolBus-Lao.pptx"

# Standard PowerPoint widescreen 16:9: 13.333" × 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # Blank

png_files = sorted(SLIDES_DIR.glob("slide-*.png"))
if not png_files:
    print(f"No slide PNGs in {SLIDES_DIR}", file=sys.stderr)
    sys.exit(1)

for png in png_files:
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(
        str(png), 0, 0, width=SLIDE_W, height=SLIDE_H
    )

prs.save(OUT)
print(f"Wrote {OUT} ({len(png_files)} slides)")
